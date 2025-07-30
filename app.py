import os
# Ensure you have the required packages installed:
import docx
import pptx
import csv
import fitz  # This is PyMuPDF
import streamlit as st
import requests # For linking connection to OpenRouter API
import tempfile # For temporary file handling and avoiding in-memory issues.
import langchain
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document

# ========== API Setup ==========
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY") or "YOUR_API_KEY"
MODEL_NAME = "deepseek/deepseek-chat-v3-0324:free"

# ========== UI Setup ==========
st.set_page_config(page_title="📄 Document Decoder", layout="wide")

st.title("📘 Welcome here")
st.markdown("""
This lets you **chat with your documents**.  
Upload one or more files, and ask anything related to their content! 💬
""")

# ========== File Upload ==========
uploaded_files = st.file_uploader("📄 Upload PDF(s)", type=["pdf","docx", "pptx", "txt", "csv"], accept_multiple_files=True)

# ========== Reset Button ==========
if st.button("🔁 Reset Session"):
    st.session_state.clear()
    st.rerun()

# ========== PDF Processing ==========
@st.cache_resource(show_spinner="📚 Indexing Document(s)... Please wait.")
def build_vector_db(uploaded_files):
    all_docs = []
    splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=50)
    # Process each uploaded file

    for file in uploaded_files:
        text = extract_text_from_file(file)
        chunks = splitter.split_text(text)
        docs = [Document(page_content=chunk, metadata={"source": file.name}) for chunk in chunks]
        all_docs.extend(docs)

    embedder = HuggingFaceEmbeddings(model_name="BAAI/bge-base-en")
    vectordb = FAISS.from_documents(all_docs, embedder)
    return vectordb.as_retriever(search_type="similarity", k=4)

# ========== Text Extraction Function ==========
@st.cache_data(show_spinner="🔍 Extracting text from files...")
def extract_text_from_file(file):
    suffix = file.name.split('.')[-1].lower()
    text = ""

    try:
        if suffix == "pdf":
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                tmp.write(file.read())
                tmp_path = tmp.name
            try:
                with fitz.open(tmp_path) as doc:
                    text = "\n".join([page.get_text() for page in doc])
            finally:
                    os.remove(tmp_path)  # ✅ Clean up temp file even if there's an error


        elif suffix == "docx":
            file.seek(0)
            doc = docx.Document(file)
            text = "\n".join([para.text for para in doc.paragraphs])

        elif suffix == "pptx":
            file.seek(0)
            presentation = pptx.Presentation(file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        elif suffix == "txt":
            file.seek(0)
            text = file.read().decode("utf-8")

        elif suffix == "csv":
            file.seek(0)
            lines = file.read().decode("utf-8").splitlines()
            reader = csv.reader(lines)
            text = "\n".join([", ".join(row) for row in reader])

        else:
            text = f"❌ Unsupported file type: {suffix}"

    except Exception as e:
        text = f"⚠️ Error while extracting {suffix}: {e}"

    return text

# ========== Chat Function ==========
def ask_deepseek(context, query):
    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "HTTP-Referer": "https://chat.openai.com",
        "X-Title": "PDF Chatbot"
    }
    messages = [
        {"role": "system", "content": "You are a helpful assistant. Use the provided context to answer questions."},
        {"role": "user", "content": f"Context:\n{context}\n\nQuestion: {query}"}
    ]
    payload = {"model": MODEL_NAME, "messages": messages}
    try:
        response = requests.post("https://openrouter.ai/api/v1/chat/completions", headers=headers, json=payload)
        data = response.json()
        return data["choices"][0]["message"]["content"]
    except Exception as e:
        return f"❌ API Error: {str(e)}"

# ========== Main Chat Flow ==========
if uploaded_files:
    st.markdown("### 🗂️ Uploaded Files")
    for file in uploaded_files:
        st.success(f"`{file.name}` loaded successfully")
        retriever = build_vector_db(uploaded_files)

    if "chat" not in st.session_state:
        st.session_state.chat = []

    query = st.text_input("💬 Ask something about the document…")

    if query:
        with st.spinner("🤖 Thinking..."):
            try:
                docs = retriever.get_relevant_documents(query)
                context = "\n\n".join([doc.page_content for doc in docs])
                answer = ask_deepseek(context, query)
            except Exception as e:
                answer = f"❌ Error: {str(e)}"
            st.session_state.chat.append({
                "question": query,
                "answer": answer,
                "sources": [doc.metadata["source"] for doc in docs]
            })

    # ========== Chat Display ==========
    for chat in reversed(st.session_state.chat):
        with st.chat_message("user"):
            st.markdown(chat["question"])
        with st.chat_message("assistant"):
            st.markdown(chat["answer"])
            for src in set(chat["sources"]):
                st.caption(f"📄 Source: `{src}`")

    # ========== Expandable Chat History ==========
    with st.expander("🕘 View Chat History"):
        for i, chat in enumerate(st.session_state.chat):
            st.markdown(f"**Q{i+1}:** {chat['question']}")
            st.markdown(f"**A{i+1}:** {chat['answer']}")
            st.markdown("---")
else:
    st.info("📥 Please upload at least one document to begin chatting.")

# ========== Footer ==========
st.markdown("""
<hr style="margin-top: 40px;">
<div style='text-align: center; color: #888; font-size: 14px;'>
    Built for Decoding.
</div>
""", unsafe_allow_html=True)
